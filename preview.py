"""Render pptx slides to PNG thumbnails by reading shapes (approximate)."""
import io, os, sys
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

prs = Presentation("EEG_Taqdimot_30_slayd.pptx")
SW = prs.slide_width
SH = prs.slide_height
SCALE = 1100 / SW
W = int(SW * SCALE)
H = int(SH * SCALE)
os.makedirs("preview", exist_ok=True)


def grad_first_color(sp):
    spPr = sp._element.spPr
    g = spPr.find(qn('a:gradFill'))
    if g is None:
        return None
    c = g.find('.//' + qn('a:srgbClr'))
    if c is not None:
        v = c.get('val')
        return tuple(int(v[i:i+2], 16) for i in (0, 2, 4))
    return None


def solid_color(sp):
    try:
        if sp.fill.type is not None and sp.fill.type == 1:  # solid
            c = sp.fill.fore_color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None


def line_color(sp):
    try:
        if sp.line.fill.type == 1:
            c = sp.line.color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None


def draw_shape(draw, img, sp):
    x0 = int(sp.left * SCALE); y0 = int(sp.top * SCALE)
    x1 = int((sp.left + sp.width) * SCALE); y1 = int((sp.top + sp.height) * SCALE)
    if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = sp.image.blob
            pim = Image.open(io.BytesIO(blob)).convert("RGBA")
            pim = pim.resize((max(1, x1 - x0), max(1, y1 - y0)), Image.LANCZOS)
            img.paste(pim, (x0, y0), pim)
        except Exception as e:
            pass
        return
    fill = grad_first_color(sp) or solid_color(sp)
    ln = line_color(sp)
    try:
        name = sp.auto_shape_type
    except Exception:
        name = None
    is_oval = (name is not None and 'OVAL' in str(name))
    is_round = (name is not None and 'ROUNDED' in str(name))
    if is_oval:
        draw.ellipse([x0, y0, x1, y1], fill=fill, outline=ln,
                     width=2 if ln else 0)
    elif is_round:
        r = max(2, int(min(x1 - x0, y1 - y0) * 0.12))
        draw.rounded_rectangle([x0, y0, x1, y1], radius=r, fill=fill,
                               outline=ln, width=2 if ln else 0)
    else:
        draw.rectangle([x0, y0, x1, y1], fill=fill, outline=ln,
                       width=2 if ln else 0)


for i, slide in enumerate(prs.slides):
    img = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    # background
    try:
        bf = slide.background.fill
        if bf.type == 1:
            c = bf.fore_color.rgb
            img.paste((c[0], c[1], c[2], 255), [0, 0, W, H])
    except Exception:
        pass
    draw = ImageDraw.Draw(img)
    for sp in slide.shapes:
        draw_shape(draw, img, sp)
    img.convert("RGB").save(f"preview/slide_{i+1:02d}.png")

# build a contact sheet of all 30
cols, rows = 5, 6
tw, th = W // 3, H // 3
sheet = Image.new("RGB", (cols * tw + (cols + 1) * 10, rows * th + (rows + 1) * 10),
                  (235, 235, 235))
for i in range(30):
    im = Image.open(f"preview/slide_{i+1:02d}.png").resize((tw, th))
    c = i % cols; r = i // cols
    sheet.paste(im, (10 + c * (tw + 10), 10 + r * (th + 10)))
sheet.save("preview/contact_sheet.png")
print("rendered", len(prs.slides._sldIdLst), "slides; contact_sheet.png ready")
