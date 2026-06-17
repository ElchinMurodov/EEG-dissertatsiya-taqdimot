"""
Builds a professional, text-free (empty placeholder) 30-slide EEG-themed
presentation in a white & blue palette, decorated with brain, EEG wave,
neuron, electrode-cap and athlete graphics.

Run: python3 build_presentation.py
Output: EEG_Taqdimot_30_slayd.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")

# ---- Palette ----
NAVY  = RGBColor(10, 38, 71)
BLUE  = RGBColor(24, 86, 160)
MID   = RGBColor(46, 111, 183)
SKY   = RGBColor(74, 144, 217)
CYAN  = RGBColor(0, 178, 216)
LIGHT = RGBColor(202, 224, 247)   # light blue
PALE  = RGBColor(238, 245, 252)   # very pale blue
PALE2 = RGBColor(228, 239, 250)
WHITE = RGBColor(255, 255, 255)

EMU = 914400
SW = int(13.333 * EMU)   # slide width  (16:9)
SH = int(7.5 * EMU)      # slide height

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def IN(v):
    return Emu(int(v * EMU))


# ---------- low-level helpers ----------
def bg(slide, color):
    el = slide.background.fill
    el.solid()
    el.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    return sp


def _soft_shadow(sp):
    spPr = sp._element.spPr
    effLst = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad': '90000', 'dist': '38100',
                           'dir': '5400000', 'rotWithShape': '0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val': '1E3A5F'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha)
    sh.append(clr)
    effLst.append(sh)
    spPr.append(effLst)


def pic_fit(slide, path, bx, by, bw, bh, align="center", valign="middle"):
    """Place image inside box keeping aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    bar = bw / bh
    if ar > bar:
        w = bw
        h = bw / ar
    else:
        h = bh
        w = bh * ar
    if align == "center":
        x = bx + (bw - w) / 2
    elif align == "left":
        x = bx
    else:
        x = bx + (bw - w)
    if valign == "middle":
        y = by + (bh - h) / 2
    elif valign == "top":
        y = by
    else:
        y = by + (bh - h)
    return slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))


def grad_band(slide, x, y, w, h, c1, c2, angle=0):
    """Rectangle with linear gradient between two colors."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    sp.line.fill.background()
    sp.shadow.inherit = False
    spPr = sp._element.spPr
    # remove default solid fill
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = spPr.makeelement(qn('a:gradFill'), {})
    lst = spPr.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = spPr.makeelement(qn('a:gs'), {'pos': str(pos)})
        c = spPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(c)
        lst.append(gs)
    grad.append(lst)
    lin = spPr.makeelement(qn('a:lin'), {'ang': str(angle * 60000), 'scaled': '1'})
    grad.append(lin)
    # insert grad fill before line element
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return sp


def dotrow(slide, x, y, n, active=0, gap=0.28, r=0.085):
    """Row of progress dots (no numbers / text)."""
    for i in range(n):
        col = CYAN if i == active else LIGHT
        rect(slide, x + i * gap, y, r, r, fill=col, shape=MSO_SHAPE.OVAL)


def corner_pulse(slide):
    """Small EEG pulse motif top-right + thin accent."""
    pic_fit(slide, os.path.join(A, "pulse.png"), 9.7, 0.28, 3.2, 0.55)


def footer(slide, idx, total):
    rect(slide, 0.9, 7.02, 11.53, 0.012, fill=LIGHT)
    pic_fit(slide, os.path.join(A, "brain.png"), 0.62, 6.74, 0.45, 0.45)
    dotrow(slide, 11.1, 6.92, min(total, 6), active=idx % 6, gap=0.26, r=0.07)


def ph(slide, x, y, w, h, fill=PALE, line=None, radius=True):
    """Empty (text-free) placeholder box."""
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = rect(slide, x, y, w, h, fill=fill, line=line,
              line_w=1.0, shape=shape)
    if radius:
        try:
            sp.adjustments[0] = 0.06
        except Exception:
            pass
    return sp


def title_ph(slide, x=0.9, y=0.62, w=8.6):
    """Header zone: accent square + thin title placeholder bar + sub bar."""
    rect(slide, x, y, 0.16, 0.66, fill=CYAN)
    ph(slide, x + 0.32, y, w, 0.5, fill=PALE2)        # title line
    ph(slide, x + 0.32, y + 0.62, w * 0.62, 0.22, fill=PALE)  # subtitle line


# ============================================================
#  SLIDE BUILDERS
# ============================================================
def s_title(slide, idx, total):
    bg(slide, WHITE)
    # left navy panel
    grad_band(slide, 0, 0, 5.6, 7.5, (10, 38, 71), (24, 86, 160), angle=90)
    # faint wave ring on panel
    pic_fit(slide, os.path.join(A, "wave_ring.png"), -1.0, 1.4, 5.0, 5.0)
    # brain (light) hero on panel
    pic_fit(slide, os.path.join(A, "brain_light.png"), 0.9, 1.9, 3.8, 3.8)
    # right white area: EEG wave strip across
    pic_fit(slide, os.path.join(A, "eeg_waves_wide.png"), 6.0, 1.05, 6.6, 1.3)
    # title placeholders (empty)
    rect(slide, 6.0, 2.75, 0.18, 1.0, fill=CYAN)
    ph(slide, 6.3, 2.78, 6.0, 0.66, fill=PALE2)       # main title
    ph(slide, 6.3, 3.62, 6.0, 0.66, fill=PALE2)       # main title 2
    ph(slide, 6.3, 4.62, 5.2, 0.3, fill=PALE)         # subtitle
    ph(slide, 6.3, 5.06, 4.4, 0.3, fill=PALE)         # subtitle 2
    # author / footer placeholders
    ph(slide, 6.3, 6.1, 2.6, 0.26, fill=PALE)
    ph(slide, 9.1, 6.1, 2.6, 0.26, fill=PALE)
    pic_fit(slide, os.path.join(A, "runner.png"), 10.7, 4.7, 1.9, 1.9)


def s_section(slide, idx, total):
    bg(slide, WHITE)
    grad_band(slide, 0, 0, 13.333, 7.5, (24, 86, 160), (10, 38, 71), angle=30)
    pic_fit(slide, os.path.join(A, "wave_ring.png"), 8.4, 0.6, 6.2, 6.2)
    # big empty number block (square outline) — text-free
    rect(slide, 1.1, 2.1, 2.4, 2.4, fill=None, line=WHITE, line_w=2.4,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, 1.42, 2.42, 1.76, 1.76, fill=RGBColor(255, 255, 255),
         shape=MSO_SHAPE.OVAL)
    pic_fit(slide, os.path.join(A, "brain.png"), 1.55, 2.55, 1.5, 1.5)
    # section title placeholders (light, on dark)
    rect(slide, 4.0, 2.55, 0.16, 1.4, fill=CYAN)
    ph(slide, 4.3, 2.6, 6.4, 0.6, fill=RGBColor(255, 255, 255))
    ph(slide, 4.3, 3.4, 5.0, 0.36, fill=RGBColor(214, 230, 247))
    ph(slide, 4.3, 3.92, 4.0, 0.36, fill=RGBColor(214, 230, 247))
    pic_fit(slide, os.path.join(A, "head_profile_light.png"), 1.0, 4.9, 1.8, 1.8)


def header(slide):
    title_ph(slide)
    corner_pulse(slide)
    # thin top wave accent line
    pic_fit(slide, os.path.join(A, "eeg_waves_wide.png"), 0.0, 0.02, 13.333, 0.4)


def s_one_big(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    ph(slide, 0.9, 1.95, 11.53, 4.7, fill=PALE)
    # a few empty text-line placeholders inside
    for i in range(6):
        ph(slide, 1.25, 2.35 + i * 0.62, 9.4 - (i % 2) * 1.2, 0.3, fill=WHITE)
    pic_fit(slide, os.path.join(A, "neuron.png"), 10.3, 4.4, 2.1, 2.1)
    footer(slide, idx, total)


def s_two_col(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    for cx in (0.9, 6.95):
        ph(slide, cx, 1.95, 5.48, 4.7, fill=PALE, line=LIGHT)
        rect(slide, cx, 1.95, 5.48, 0.5, fill=MID,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        for i in range(5):
            ph(slide, cx + 0.3, 2.7 + i * 0.66, 4.7 - (i % 2) * 0.8, 0.3, fill=WHITE)
    # small icons in column headers
    pic_fit(slide, os.path.join(A, "brain.png"), 5.9, 2.0, 0.4, 0.4)
    pic_fit(slide, os.path.join(A, "electrode_cap.png"), 11.95, 2.0, 0.4, 0.4)
    footer(slide, idx, total)


def s_content_image(slide, idx, total, icon, on_left=False):
    bg(slide, WHITE)
    header(slide)
    if on_left:
        # image area left
        rect(slide, 0.9, 1.95, 4.9, 4.7, fill=PALE2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LIGHT)
        pic_fit(slide, os.path.join(A, icon), 1.2, 2.25, 4.3, 4.1)
        tx = 6.1
    else:
        rect(slide, 7.55, 1.95, 4.88, 4.7, fill=PALE2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LIGHT)
        pic_fit(slide, os.path.join(A, icon), 7.8, 2.25, 4.4, 4.1)
        tx = 0.9
    ph(slide, tx, 1.95, 5.3, 0.5, fill=PALE2)
    for i in range(6):
        ph(slide, tx, 2.7 + i * 0.66, 5.3 - (i % 2) * 1.0, 0.3, fill=PALE)
    footer(slide, idx, total)


def s_three_cards(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    icons = ["brain.png", "electrode_cap.png", "runner.png"]
    xs = [0.9, 5.12, 9.34]
    for i, x in enumerate(xs):
        rect(slide, x, 1.95, 3.6, 4.7, fill=PALE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LIGHT, shadow=True)
        # icon circle
        rect(slide, x + 1.2, 2.25, 1.2, 1.2, fill=WHITE, shape=MSO_SHAPE.OVAL)
        pic_fit(slide, os.path.join(A, icons[i]), x + 1.32, 2.37, 0.96, 0.96)
        ph(slide, x + 0.4, 3.7, 2.8, 0.36, fill=PALE2)   # card title
        for j in range(4):
            ph(slide, x + 0.4, 4.25 + j * 0.5, 2.8 - (j % 2) * 0.5, 0.26, fill=WHITE)
    footer(slide, idx, total)


def s_stat_circles(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    xs = [1.1, 4.18, 7.26, 10.34]
    for i, x in enumerate(xs):
        rect(slide, x, 2.4, 1.9, 1.9, fill=None, line=SKY, line_w=2.2,
             shape=MSO_SHAPE.OVAL)
        rect(slide, x + 0.28, 2.68, 1.34, 1.34, fill=PALE, shape=MSO_SHAPE.OVAL)
        ph(slide, x + 0.45, 3.18, 1.0, 0.34, fill=WHITE)   # big number placeholder
        ph(slide, x + 0.1, 4.55, 1.7, 0.28, fill=PALE)     # label
        ph(slide, x + 0.25, 4.92, 1.4, 0.24, fill=PALE)
    pic_fit(slide, os.path.join(A, "eeg_waves_wide.png"), 0.9, 5.7, 11.53, 0.8)
    footer(slide, idx, total)


def s_full_eeg(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    rect(slide, 0.9, 1.95, 11.53, 3.5, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pic_fit(slide, os.path.join(A, "eeg_waves_block.png"), 1.15, 2.15, 11.0, 3.1)
    # caption placeholders below
    for i in range(2):
        ph(slide, 0.9 + i * 5.9, 5.7, 5.5, 0.3, fill=PALE)
    for i in range(2):
        ph(slide, 0.9 + i * 5.9, 6.12, 4.4, 0.26, fill=PALE)
    footer(slide, idx, total)


def s_quote(slide, idx, total):
    bg(slide, WHITE)
    grad_band(slide, 0, 0, 13.333, 7.5, (238, 245, 252), (255, 255, 255), angle=60)
    rect(slide, 0.0, 0.0, 0.22, 7.5, fill=CYAN)
    pic_fit(slide, os.path.join(A, "head_profile.png"), 9.7, 1.5, 3.4, 3.4)
    pic_fit(slide, os.path.join(A, "pulse.png"), 1.2, 1.4, 6.0, 0.7)
    # big quote placeholder lines
    for i in range(3):
        ph(slide, 1.2, 2.7 + i * 0.8, 8.0 - i * 0.8, 0.5, fill=PALE2)
    ph(slide, 1.2, 5.5, 3.0, 0.3, fill=PALE)   # attribution
    footer(slide, idx, total)


def s_timeline(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    # horizontal line
    rect(slide, 1.1, 4.0, 11.1, 0.06, fill=LIGHT)
    xs = [1.4, 3.6, 5.8, 8.0, 10.2]
    for i, x in enumerate(xs):
        rect(slide, x, 3.78, 0.5, 0.5, fill=CYAN if i % 2 == 0 else MID,
             shape=MSO_SHAPE.OVAL)
        above = (i % 2 == 0)
        by = 2.2 if above else 4.7
        ph(slide, x - 0.8, by, 2.1, 0.34, fill=PALE2)
        for j in range(2):
            ph(slide, x - 0.8, by + 0.45 + j * 0.34, 2.1 - j * 0.5, 0.24, fill=PALE)
    footer(slide, idx, total)


def s_compare(slide, idx, total):
    bg(slide, WHITE)
    header(slide)
    panels = [(0.9, NAVY, "brain.png"), (6.95, MID, "runner.png")]
    for x, col, icon in panels:
        rect(slide, x, 1.95, 5.48, 4.7, fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(slide, x + 0.3, 2.25, 1.0, 1.0, fill=WHITE, shape=MSO_SHAPE.OVAL)
        pic_fit(slide, os.path.join(A, icon), x + 0.4, 2.35, 0.8, 0.8)
        ph(slide, x + 1.5, 2.45, 3.5, 0.5, fill=RGBColor(214, 230, 247))
        for j in range(5):
            ph(slide, x + 0.4, 3.6 + j * 0.58, 4.6 - (j % 2) * 0.8, 0.3,
               fill=RGBColor(235, 243, 251))
    # center VS circle
    rect(slide, 6.27, 3.9, 0.78, 0.78, fill=CYAN, shape=MSO_SHAPE.OVAL, shadow=True)
    footer(slide, idx, total)


def s_closing(slide, idx, total):
    bg(slide, WHITE)
    grad_band(slide, 0, 0, 13.333, 7.5, (10, 38, 71), (24, 86, 160), angle=120)
    pic_fit(slide, os.path.join(A, "wave_ring.png"), 8.0, -0.6, 6.6, 6.6)
    pic_fit(slide, os.path.join(A, "brain_light.png"), 0.8, 1.0, 3.4, 3.4)
    pic_fit(slide, os.path.join(A, "eeg_waves_wide.png"), 0.7, 4.7, 8.0, 1.1)
    rect(slide, 4.6, 1.4, 0.16, 1.5, fill=CYAN)
    ph(slide, 4.9, 1.45, 5.6, 0.7, fill=RGBColor(255, 255, 255))   # "Thank you"
    ph(slide, 4.9, 2.4, 4.4, 0.4, fill=RGBColor(214, 230, 247))
    # contact placeholders
    for i in range(3):
        ph(slide, 4.9, 3.4 + i * 0.5, 4.0, 0.3, fill=RGBColor(214, 230, 247))
    pic_fit(slide, os.path.join(A, "runner_light.png"), 11.0, 5.0, 1.9, 1.9)


# ============================================================
#  ASSEMBLE 30 SLIDES
# ============================================================
TOTAL = 30
plan = [
    ("title", {}),                                   # 1
    ("two_col", {}),                                 # 2  overview/agenda
    ("section", {}),                                 # 3  divider 01
    ("one_big", {}),                                 # 4
    ("content_image", {"icon": "head_profile.png"}), # 5
    ("two_col", {}),                                 # 6
    ("three_cards", {}),                             # 7
    ("full_eeg", {}),                                # 8
    ("section", {}),                                 # 9  divider 02
    ("content_image", {"icon": "electrode_cap.png", "on_left": True}),  # 10
    ("stat_circles", {}),                            # 11
    ("two_col", {}),                                 # 12
    ("one_big", {}),                                 # 13
    ("section", {}),                                 # 14 divider 03
    ("content_image", {"icon": "runner.png", "on_left": True}),  # 15
    ("three_cards", {}),                             # 16
    ("full_eeg", {}),                                # 17
    ("compare", {}),                                 # 18
    ("one_big", {}),                                 # 19
    ("section", {}),                                 # 20 divider 04
    ("content_image", {"icon": "brain.png"}),        # 21
    ("compare", {}),                                 # 22
    ("timeline", {}),                                # 23
    ("three_cards", {}),                             # 24
    ("section", {}),                                 # 25 divider 05
    ("content_image", {"icon": "neuron.png", "on_left": True}),  # 26
    ("two_col", {}),                                 # 27
    ("quote", {}),                                   # 28
    ("stat_circles", {}),                            # 29
    ("closing", {}),                                 # 30
]

builders = {
    "title": s_title,
    "section": s_section,
    "one_big": s_one_big,
    "two_col": s_two_col,
    "content_image": s_content_image,
    "three_cards": s_three_cards,
    "stat_circles": s_stat_circles,
    "full_eeg": s_full_eeg,
    "quote": s_quote,
    "timeline": s_timeline,
    "compare": s_compare,
    "closing": s_closing,
}

for i, (kind, kw) in enumerate(plan):
    slide = prs.slides.add_slide(BLANK)
    builders[kind](slide, i, TOTAL, **kw)

OUT = os.path.join(HERE, "EEG_Taqdimot_30_slayd.pptx")
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
